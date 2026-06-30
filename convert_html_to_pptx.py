import os
import sys
import time
import socket
import threading
import http.server
import socketserver
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.options import Options as ChromeOptions
from selenium.webdriver.edge.options import Options as EdgeOptions
from pptx import Presentation
from pptx.util import Inches

# Find a free port
def get_free_port():
    s = socket.socket()
    s.bind(('', 0))
    port = s.getsockname()[1]
    s.close()
    return port

# Simple HTTP server thread
class HTTPServerThread(threading.Thread):
    def __init__(self, port, directory):
        super().__init__()
        self.port = port
        self.directory = directory
        self.httpd = None

    def run(self):
        handler = lambda *args, **kwargs: http.server.SimpleHTTPRequestHandler(*args, directory=self.directory, **kwargs)
        class SilentTCPServer(socketserver.TCPServer):
            allow_reuse_address = True
        
        self.httpd = SilentTCPServer(("", self.port), handler)
        self.httpd.serve_forever()

    def stop(self):
        if self.httpd:
            self.httpd.shutdown()
            self.httpd.server_close()

def main():
    workspace = r"c:\Users\ACER\Documents\PREPARE TA\ppt"
    port = get_free_port()
    
    # Start HTTP server
    print(f"Starting temporary local server on port {port}...")
    server_thread = HTTPServerThread(port, workspace)
    server_thread.daemon = True
    server_thread.start()
    time.sleep(1.5) # Let server start
    
    # Setup selenium driver (Try Chrome, then Edge)
    driver = None
    
    chrome_opts = ChromeOptions()
    chrome_opts.add_argument("--headless=new")
    chrome_opts.add_argument("--window-size=1400,900")
    chrome_opts.add_argument("--disable-gpu")
    chrome_opts.add_argument("--no-sandbox")
    chrome_opts.add_argument("--disable-dev-shm-usage")
    
    print("Initializing browser automation...")
    try:
        driver = webdriver.Chrome(options=chrome_opts)
        print("Successfully opened Chrome via Selenium.")
    except Exception as e:
        print(f"Chrome failed: {e}. Trying Edge...")
        try:
            edge_opts = EdgeOptions()
            edge_opts.add_argument("--headless=new")
            edge_opts.add_argument("--window-size=1400,900")
            edge_opts.add_argument("--disable-gpu")
            edge_opts.add_argument("--no-sandbox")
            driver = webdriver.Edge(options=edge_opts)
            print("Successfully opened Edge via Selenium.")
        except Exception as e2:
            print(f"Edge failed: {e2}.")
            server_thread.stop()
            sys.exit("Error: Could not start Chrome or Edge webdriver. Please make sure Google Chrome or Microsoft Edge is installed.")

    try:
        # Load the presentation page
        url = f"http://localhost:{port}/ppt.html"
        print(f"Loading presentation page: {url}")
        driver.get(url)
        time.sleep(2.5) # Wait for page load and content.js to execute
        
        # Inject style to disable transitions/animations for instant screenshots
        print("Disabling transitions and scaling stage to 100%...")
        disable_transitions_js = """
        const style = document.createElement('style');
        style.innerHTML = `
            * {
                transition: none !important;
                animation: none !important;
                transition-duration: 0s !important;
                animation-duration: 0s !important;
            }
        `;
        document.head.appendChild(style);
        
        // Force stage scale to 1.0 at top-left, hide scrollbars
        document.body.style.overflow = 'hidden';
        const stage = document.getElementById('stage');
        stage.style.transform = 'scale(1)';
        stage.style.transformOrigin = 'top left';
        stage.style.position = 'absolute';
        stage.style.top = '0';
        stage.style.left = '0';
        """
        driver.execute_script(disable_transitions_js)
        time.sleep(0.5)
        
        # Get total slides from the page
        TOTAL_SLIDES = int(driver.execute_script("return TOTAL;"))
        print(f"Detected {TOTAL_SLIDES} slides in total.")
        
        temp_dir = os.path.join(workspace, "temp_slides")
        os.makedirs(temp_dir, exist_ok=True)
        
        image_paths = []
        
        # Capture screenshots
        for i in range(1, TOTAL_SLIDES + 1):
            print(f"Capturing Slide {i}/{TOTAL_SLIDES}...")
            driver.execute_script(f"showSlide({i});")
            time.sleep(0.25) # short buffer for image rendering if any
            
            # Capture stage screenshot
            stage_element = driver.find_element(By.ID, "stage")
            img_path = os.path.join(temp_dir, f"slide_{i}.png")
            stage_element.screenshot(img_path)
            image_paths.append(img_path)
            
        print("All slides captured successfully. Building PowerPoint file...")
        
        # Create PowerPoint
        prs = Presentation()
        # Set to 16:9 widescreen ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]
        
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
        output_pptx = os.path.join(workspace, "Presentasi_PLTS_Hybrid_IoT.pptx")
        try:
            prs.save(output_pptx)
            print(f"PowerPoint presentation saved successfully to: {output_pptx}")
        except PermissionError:
            alternative_path = os.path.join(workspace, "Presentasi_PLTS_Hybrid_IoT_new.pptx")
            print(f"Warning: Permission denied for '{output_pptx}' (it might be open in PowerPoint).")
            print(f"Saving to alternative path: '{alternative_path}'")
            prs.save(alternative_path)
            print(f"PowerPoint presentation saved successfully to: {alternative_path}")
        
        # Cleanup temporary slide images
        print("Cleaning up temporary images...")
        for img_path in image_paths:
            try:
                os.remove(img_path)
            except:
                pass
        try:
            os.rmdir(temp_dir)
        except:
            pass
            
        print("Conversion process complete!")
        
    finally:
        driver.quit()
        server_thread.stop()

if __name__ == "__main__":
    main()
